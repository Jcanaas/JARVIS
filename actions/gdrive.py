"""Google Drive integration for Jarvis."""
from __future__ import annotations

import io
import mimetypes
import zipfile
from xml.etree import ElementTree
from pathlib import Path
from typing import Dict, List, Optional


GOOGLE_EXPORTS = {
    "application/vnd.google-apps.document": ("application/pdf", ".pdf"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    ),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    ),
    "application/vnd.google-apps.drawing": ("image/png", ".png"),
}


def _get_service():
    from actions.google_auth import get_google_service
    return get_google_service("drive", "v3")


def _emit_progress(progress_hook, *, active: bool, percent: float, label: str, detail: str, can_cancel: bool = False):
    if not progress_hook:
        return
    try:
        progress_hook({
            "title": "GOOGLE DRIVE",
            "active": active,
            "percent": percent,
            "label": label,
            "detail": detail,
            "can_cancel": can_cancel,
        })
    except Exception:
        pass


def _escape_query(value: str) -> str:
    return str(value or "").replace("\\", "\\\\").replace("'", "\\'")


def _fmt_size(size_str: str) -> str:
    try:
        b = int(size_str)
        if b >= 1_073_741_824:
            return f"{b / 1_073_741_824:.1f} GB"
        if b >= 1_048_576:
            return f"{b / 1_048_576:.1f} MB"
        if b >= 1024:
            return f"{b / 1024:.1f} KB"
        return f"{b} B"
    except Exception:
        return ""


def _as_bool(value) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on", "si", "sí"}


def _downloads_dir() -> Path:
    out = Path.home() / "Downloads" / "JARVIS_Drive"
    out.mkdir(parents=True, exist_ok=True)
    return out


def _find_folder(svc, name: str, parent_id: Optional[str] = None) -> Optional[str]:
    safe_name = _escape_query(name)
    q = f"name='{safe_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
    if parent_id:
        q += f" and '{parent_id}' in parents"
    resp = svc.files().list(q=q, fields="files(id,name)", spaces="drive").execute()
    files = resp.get("files", [])
    return files[0]["id"] if files else None


def _find_or_create_folder(svc, name: str, parent_id: Optional[str] = None) -> str:
    fid = _find_folder(svc, name, parent_id)
    if fid:
        return fid
    meta = {"name": name, "mimeType": "application/vnd.google-apps.folder"}
    if parent_id:
        meta["parents"] = [parent_id]
    f = svc.files().create(body=meta, fields="id").execute()
    return f["id"]


def _file_fields() -> str:
    return (
        "id,name,mimeType,size,modifiedTime,createdTime,description,"
        "webViewLink,webContentLink,thumbnailLink,iconLink,parents,trashed"
    )


def get_file_info(file_id: str) -> Dict:
    svc = _get_service()
    return svc.files().get(fileId=file_id, fields=_file_fields()).execute()


def _office_text_preview(data: bytes, mime: str, name: str) -> Optional[str]:
    """Extract readable text from common Office formats without opening a browser."""
    suffix = Path(name).suffix.lower()
    try:
        if mime.endswith("wordprocessingml.document") or suffix == ".docx":
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                root = ElementTree.fromstring(archive.read("word/document.xml"))
            paragraphs = []
            for paragraph in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p"):
                text = "".join(
                    node.text or ""
                    for node in paragraph.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t")
                ).strip()
                if text:
                    paragraphs.append(text)
            return "\n\n".join(paragraphs)

        if mime.endswith("presentationml.presentation") or suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(io.BytesIO(data))
            slides = []
            for index, slide in enumerate(presentation.slides, 1):
                lines = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if lines:
                    slides.append(f"DIAPOSITIVA {index}\n" + "\n".join(lines))
            return "\n\n".join(slides)

        if mime.endswith("spreadsheetml.sheet") or suffix == ".xlsx":
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                shared = []
                if "xl/sharedStrings.xml" in archive.namelist():
                    root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
                    ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
                    shared = [
                        "".join(node.text or "" for node in item.iter(f"{ns}t"))
                        for item in root.iter(f"{ns}si")
                    ]
                sheet_names = sorted(
                    path for path in archive.namelist()
                    if path.startswith("xl/worksheets/sheet") and path.endswith(".xml")
                )
                rows = []
                ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
                for sheet_index, sheet_name in enumerate(sheet_names[:5], 1):
                    rows.append(f"HOJA {sheet_index}")
                    root = ElementTree.fromstring(archive.read(sheet_name))
                    for row in list(root.iter(f"{ns}row"))[:200]:
                        values = []
                        for cell in row.iter(f"{ns}c"):
                            value = cell.find(f"{ns}v")
                            raw = value.text if value is not None else ""
                            if cell.get("t") == "s" and raw.isdigit():
                                position = int(raw)
                                raw = shared[position] if position < len(shared) else raw
                            values.append(raw)
                        rows.append("\t".join(values))
                    rows.append("")
                return "\n".join(rows)
    except Exception:
        return None
    return None


def get_file_preview(file_id: str, max_bytes: int = 64 * 1024 * 1024) -> Dict:
    """Download a bounded representation suitable for an in-app preview."""
    from googleapiclient.http import MediaIoBaseDownload

    if not file_id:
        raise ValueError("Necesito el ID del archivo.")
    svc = _get_service()
    info = get_file_info(file_id)
    mime = str(info.get("mimeType") or "")
    name = str(info.get("name") or file_id)

    if mime == "application/vnd.google-apps.folder":
        return {"kind": "folder", "info": info, "data": b""}

    if mime.startswith("video/") and info.get("thumbnailLink"):
        try:
            from google.auth.transport.requests import AuthorizedSession

            credentials = getattr(getattr(svc, "_http", None), "credentials", None)
            if credentials is not None:
                response = AuthorizedSession(credentials).get(
                    info["thumbnailLink"],
                    timeout=15,
                )
                response.raise_for_status()
                return {
                    "kind": "image",
                    "info": info,
                    "data": response.content,
                    "mimeType": response.headers.get("content-type", "image/jpeg"),
                    "name": name,
                    "video": True,
                }
        except Exception:
            pass

    if mime in {
        "application/vnd.google-apps.document",
        "application/vnd.google-apps.presentation",
        "application/vnd.google-apps.drawing",
    }:
        request = svc.files().export_media(fileId=file_id, mimeType="application/pdf")
        preview_mime = "application/pdf"
    elif mime == "application/vnd.google-apps.spreadsheet":
        request = svc.files().export_media(fileId=file_id, mimeType="text/csv")
        preview_mime = "text/csv"
    else:
        size = int(info.get("size") or 0)
        if size and size > max_bytes:
            return {"kind": "too_large", "info": info, "data": b"", "size": size}
        request = svc.files().get_media(fileId=file_id)
        preview_mime = mime

    output = io.BytesIO()
    downloader = MediaIoBaseDownload(output, request, chunksize=min(max_bytes, 1024 * 1024))
    done = False
    while not done:
        _status, done = downloader.next_chunk()
        if output.tell() > max_bytes:
            return {
                "kind": "too_large",
                "info": info,
                "data": b"",
                "size": output.tell(),
            }

    data = output.getvalue()
    if preview_mime.startswith("image/"):
        kind = "image"
    elif preview_mime.startswith("audio/"):
        kind = "audio"
    elif preview_mime == "application/pdf":
        kind = "pdf"
    elif preview_mime.startswith("text/") or preview_mime in {
        "application/json",
        "application/xml",
        "application/javascript",
    }:
        kind = "text"
    else:
        office_text = _office_text_preview(data, preview_mime, name)
        if office_text is not None:
            kind = "text"
            data = office_text.encode("utf-8")
            preview_mime = "text/plain"
        else:
            kind = "unsupported"
    return {
        "kind": kind,
        "info": info,
        "data": data,
        "mimeType": preview_mime,
        "name": name,
    }


def list_files(count: int = 20, folder_id: Optional[str] = None) -> List[Dict]:
    svc = _get_service()
    q = "trashed=false"
    if folder_id:
        q += f" and '{folder_id}' in parents"
    resp = svc.files().list(
        q=q,
        pageSize=count,
        fields=f"files({_file_fields()})",
        orderBy="modifiedTime desc",
        spaces="drive",
    ).execute()
    return resp.get("files", [])


def search_files(
    query: str,
    count: int = 20,
    mime_type: str = "",
    folder_id: Optional[str] = None,
) -> List[Dict]:
    svc = _get_service()
    safe = _escape_query(query)
    q = f"name contains '{safe}' and trashed=false"
    if mime_type:
        q += f" and mimeType='{_escape_query(mime_type)}'"
    if folder_id:
        q += f" and '{_escape_query(folder_id)}' in parents"
    resp = svc.files().list(
        q=q,
        pageSize=count,
        fields=f"files({_file_fields()})",
        orderBy="modifiedTime desc",
        spaces="drive",
    ).execute()
    return resp.get("files", [])


def _resolve_file_id(svc, file_id: str = "", query: str = "") -> Optional[str]:
    if file_id:
        return file_id
    if not query:
        return None
    safe = _escape_query(query)
    resp = svc.files().list(
        q=f"name contains '{safe}' and trashed=false",
        pageSize=1,
        fields="files(id)",
        orderBy="modifiedTime desc",
        spaces="drive",
    ).execute()
    files = resp.get("files", [])
    return files[0]["id"] if files else None


def upload_file(
    file_path: str,
    folder_name: Optional[str] = None,
    folder_id: Optional[str] = None,
    name: str = "",
    progress_hook=None,
) -> Dict:
    from googleapiclient.http import MediaFileUpload

    svc = _get_service()
    path = Path(file_path).expanduser()
    if not path.exists() or not path.is_file():
        raise FileNotFoundError(f"Archivo no encontrado: {file_path}")

    target_folder_id = folder_id
    if folder_name and not target_folder_id:
        target_folder_id = _find_or_create_folder(svc, folder_name)

    mime, _ = mimetypes.guess_type(str(path))
    mime = mime or "application/octet-stream"

    file_meta = {"name": name or path.name}
    if target_folder_id:
        file_meta["parents"] = [target_folder_id]

    _emit_progress(progress_hook, active=True, percent=0, label="Subiendo", detail=path.name)
    media = MediaFileUpload(str(path), mimetype=mime, resumable=True)
    request = svc.files().create(body=file_meta, media_body=media, fields=_file_fields())
    response = None
    while response is None:
        status, response = request.next_chunk()
        if status:
            _emit_progress(
                progress_hook,
                active=True,
                percent=float(status.progress() * 100),
                label="Subiendo",
                detail=path.name,
            )
    _emit_progress(progress_hook, active=False, percent=100, label="Subida completa", detail=response.get("name", path.name))
    return response


def download_file(
    file_id: str = "",
    query: str = "",
    output_dir: str = "",
    export_mime: str = "",
    progress_hook=None,
) -> Dict:
    from googleapiclient.http import MediaIoBaseDownload

    svc = _get_service()
    resolved_id = _resolve_file_id(svc, file_id=file_id, query=query)
    if not resolved_id:
        raise FileNotFoundError("No se encontro el archivo en Drive.")

    info = get_file_info(resolved_id)
    outdir = Path(output_dir).expanduser() if output_dir else _downloads_dir()
    outdir.mkdir(parents=True, exist_ok=True)

    mime = info.get("mimeType", "")
    name = info.get("name", resolved_id)
    if mime in GOOGLE_EXPORTS:
        default_export, ext = GOOGLE_EXPORTS[mime]
        chosen_mime = export_mime or default_export
        request = svc.files().export_media(fileId=resolved_id, mimeType=chosen_mime)
        if not Path(name).suffix:
            name = f"{name}{ext}"
    else:
        request = svc.files().get_media(fileId=resolved_id)

    target = outdir / name
    _emit_progress(progress_hook, active=True, percent=0, label="Descargando", detail=name)
    with target.open("wb") as fh:
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            if status:
                _emit_progress(
                    progress_hook,
                    active=True,
                    percent=float(status.progress() * 100),
                    label="Descargando",
                    detail=name,
                )
    _emit_progress(progress_hook, active=False, percent=100, label="Descarga completa", detail=str(target))
    return {"id": resolved_id, "name": name, "path": str(target), "mimeType": mime}


def create_folder(name: str, parent_id: Optional[str] = None) -> Dict:
    svc = _get_service()
    meta = {"name": name, "mimeType": "application/vnd.google-apps.folder"}
    if parent_id:
        meta["parents"] = [parent_id]
    return svc.files().create(body=meta, fields="id,name,webViewLink").execute()


def share_file(
    file_id: str = "",
    query: str = "",
    email: str = "",
    role: str = "reader",
    anyone: bool = False,
    notify: bool = False,
) -> Dict:
    svc = _get_service()
    resolved_id = _resolve_file_id(svc, file_id=file_id, query=query)
    if not resolved_id:
        raise FileNotFoundError("No se encontro el archivo para compartir.")

    role = role if role in {"reader", "commenter", "writer"} else "reader"
    if anyone:
        body = {"type": "anyone", "role": role}
    else:
        if not email:
            raise ValueError("Necesito el email para compartir, o anyone=true para link publico.")
        body = {"type": "user", "role": role, "emailAddress": email}

    perm = svc.permissions().create(
        fileId=resolved_id,
        body=body,
        sendNotificationEmail=bool(notify),
        fields="id,type,role,emailAddress",
    ).execute()
    info = get_file_info(resolved_id)
    return {"permission": perm, "file": info}


def rename_file(file_id: str = "", query: str = "", new_name: str = "") -> Dict:
    svc = _get_service()
    resolved_id = _resolve_file_id(svc, file_id=file_id, query=query)
    if not resolved_id:
        raise FileNotFoundError("No se encontro el archivo para renombrar.")
    if not new_name:
        raise ValueError("Necesito el nuevo nombre.")
    return svc.files().update(fileId=resolved_id, body={"name": new_name}, fields=_file_fields()).execute()


def update_file(
    file_id: str = "",
    query: str = "",
    file_path: str = "",
    new_name: str = "",
    description: str = "",
    progress_hook=None,
) -> Dict:
    from googleapiclient.http import MediaFileUpload

    svc = _get_service()
    resolved_id = _resolve_file_id(svc, file_id=file_id, query=query)
    if not resolved_id:
        raise FileNotFoundError("No se encontro el archivo para actualizar.")

    body = {}
    if new_name:
        body["name"] = new_name
    if description:
        body["description"] = description

    media = None
    detail = new_name or resolved_id
    if file_path:
        path = Path(file_path).expanduser()
        if not path.exists() or not path.is_file():
            raise FileNotFoundError(f"Archivo local no encontrado: {file_path}")
        mime, _ = mimetypes.guess_type(str(path))
        media = MediaFileUpload(str(path), mimetype=mime or "application/octet-stream", resumable=True)
        detail = path.name

    _emit_progress(progress_hook, active=True, percent=0, label="Actualizando", detail=detail)
    request = svc.files().update(fileId=resolved_id, body=body, media_body=media, fields=_file_fields())
    if media is None:
        result = request.execute()
    else:
        result = None
        while result is None:
            status, result = request.next_chunk()
            if status:
                _emit_progress(
                    progress_hook,
                    active=True,
                    percent=float(status.progress() * 100),
                    label="Actualizando",
                    detail=detail,
                )
    _emit_progress(progress_hook, active=False, percent=100, label="Actualizacion completa", detail=result.get("name", detail))
    return result


def delete_file(file_id: str = "", query: str = "", permanent: bool = False) -> str:
    svc = _get_service()
    resolved_id = _resolve_file_id(svc, file_id=file_id, query=query)
    if not resolved_id:
        raise FileNotFoundError("No se encontro el archivo para borrar.")
    if permanent:
        svc.files().delete(fileId=resolved_id).execute()
        return f"Archivo eliminado definitivamente: {resolved_id}"
    svc.files().update(fileId=resolved_id, body={"trashed": True}, fields="id,trashed").execute()
    return f"Archivo movido a la papelera: {resolved_id}"


def _format_files(title: str, files: List[Dict]) -> str:
    if not files:
        return "No se encontraron archivos en Drive."
    lines = [title]
    for f in files:
        size_str = f" ({_fmt_size(f.get('size', ''))})" if f.get("size") else ""
        date_str = (f.get("modifiedTime") or "")[:10]
        link = f.get("webViewLink", "")
        lines.append(f"- {f.get('name', '')}{size_str} | {date_str} | ID: {f.get('id', '')} | {link}")
    return "\n".join(lines)


def gdrive(parameters: dict, player=None, speak=None, progress_hook=None) -> str:
    action = str(parameters.get("action", "list_files")).lower().strip()

    try:
        if action == "list_files":
            count = int(parameters.get("count") or parameters.get("limit") or 20)
            files = list_files(count=count, folder_id=parameters.get("folder_id"))
            return _format_files(f"{len(files)} archivo(s) recientes en Drive:", files)

        if action == "search_files":
            query = parameters.get("query") or parameters.get("name") or ""
            if not query:
                return "Necesito un termino de busqueda."
            count = int(parameters.get("count") or parameters.get("limit") or 20)
            files = search_files(query=query, count=count, mime_type=parameters.get("mime_type") or "")
            return _format_files(f"Resultados para '{query}':", files)

        if action == "upload_file":
            file_path = parameters.get("file_path") or parameters.get("path") or ""
            if not file_path:
                return "Necesito la ruta del archivo local."
            f = upload_file(
                file_path=file_path,
                folder_name=parameters.get("folder_name") or parameters.get("folder"),
                folder_id=parameters.get("folder_id"),
                name=parameters.get("name") or parameters.get("new_name") or "",
                progress_hook=progress_hook,
            )
            return f"'{f.get('name')}' subido a Drive. ID: {f.get('id')} Link: {f.get('webViewLink', '')}"

        if action == "download_file":
            f = download_file(
                file_id=parameters.get("file_id") or parameters.get("id") or "",
                query=parameters.get("query") or parameters.get("name") or "",
                output_dir=parameters.get("output_dir") or parameters.get("path") or "",
                export_mime=parameters.get("export_mime") or "",
                progress_hook=progress_hook,
            )
            return f"Archivo descargado: {f['name']} -> {f['path']}"

        if action == "create_folder":
            name = parameters.get("folder_name") or parameters.get("name") or ""
            if not name:
                return "Necesito el nombre de la carpeta."
            f = create_folder(name, parameters.get("folder_id") or parameters.get("parent_id"))
            return f"Carpeta '{f['name']}' creada en Drive. ID: {f['id']}"

        if action == "share_file":
            f = share_file(
                file_id=parameters.get("file_id") or parameters.get("id") or "",
                query=parameters.get("query") or parameters.get("name") or "",
                email=parameters.get("email") or parameters.get("to") or "",
                role=parameters.get("role") or "reader",
                anyone=_as_bool(parameters.get("anyone", False)),
                notify=_as_bool(parameters.get("notify", False)),
            )
            link = f.get("file", {}).get("webViewLink", "")
            role = f.get("permission", {}).get("role", "reader")
            return f"Archivo compartido con permiso {role}. Link: {link}"

        if action in ("rename_file", "rename"):
            f = rename_file(
                file_id=parameters.get("file_id") or parameters.get("id") or "",
                query=parameters.get("query") or parameters.get("name") or "",
                new_name=parameters.get("new_name") or parameters.get("title") or "",
            )
            return f"Archivo renombrado: {f.get('name')} | ID: {f.get('id')}"

        if action in ("update_file", "replace_file"):
            f = update_file(
                file_id=parameters.get("file_id") or parameters.get("id") or "",
                query=parameters.get("query") or parameters.get("name") or "",
                file_path=parameters.get("file_path") or parameters.get("local_path") or "",
                new_name=parameters.get("new_name") or "",
                description=parameters.get("description") or "",
                progress_hook=progress_hook,
            )
            return f"Archivo actualizado: {f.get('name')} | ID: {f.get('id')}"

        if action in ("delete_file", "trash_file"):
            return delete_file(
                file_id=parameters.get("file_id") or parameters.get("id") or "",
                query=parameters.get("query") or parameters.get("name") or "",
                permanent=_as_bool(parameters.get("permanent", False)),
            )

        if action == "get_file_info":
            file_id = parameters.get("file_id") or parameters.get("id") or ""
            if not file_id:
                return "Necesito el ID del archivo."
            return str(get_file_info(file_id))

        return (
            f"Accion desconocida: {action}. Usa list_files, search_files, upload_file, "
            "download_file, create_folder, share_file, rename_file, update_file, delete_file o get_file_info."
        )

    except FileNotFoundError as e:
        return str(e)
    except Exception as e:
        err = str(e)
        if any(k in err.lower() for k in ("invalid_grant", "token", "credentials", "unauthorized", "403", "401")):
            try:
                from actions.auth_dialog import show_gcal_setup_dialog
                show_gcal_setup_dialog()
            except Exception:
                pass
            return f"Google Drive: error de autenticacion - {err}"
        return f"Google Drive error: {err}"
